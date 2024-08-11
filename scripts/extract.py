import os
import sys
import magic
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import csv
import xlrd
import openpyxl

def clean_text(text):
    # Remove newlines and extra spaces
    return ' '.join(text.split())

def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return clean_text(text)

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + " "
    return clean_text(text)

def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return clean_text(text)

def extract_text_from_csv(file_path):
    text = ""
    with open(file_path, newline='', encoding='utf-8') as csvfile:
        reader = csv.reader(csvfile)
        for row in reader:
            text += ', '.join(row) + " "
    return clean_text(text)

def extract_text_from_xlsx(file_path):
    workbook = openpyxl.load_workbook(file_path, data_only=True)
    text = ""
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            text += ', '.join([str(cell) if cell is not None else '' for cell in row]) + " "
    return clean_text(text)

def extract_text_from_xls(file_path):
    workbook = xlrd.open_workbook(file_path)
    text = ""
    for sheet in workbook.sheets():
        for row_idx in range(sheet.nrows):
            row = sheet.row(row_idx)
            text += ', '.join([str(cell.value) for cell in row]) + " "
    return clean_text(text)

def extract_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return clean_text(file.read())

def detect_and_extract_text(file_path):
    mime = magic.Magic(mime=True)
    file_type = mime.from_file(file_path)

    if file_type == 'application/pdf':
        return extract_text_from_pdf(file_path)
    elif file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        return extract_text_from_docx(file_path)
    elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        return extract_text_from_pptx(file_path)
    elif file_type == 'text/csv':
        return extract_text_from_csv(file_path)
    elif file_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
        return extract_text_from_xlsx(file_path)
    elif file_type == 'application/vnd.ms-excel':
        return extract_text_from_xls(file_path)
    elif file_type == 'text/plain':
        return extract_text_from_txt(file_path)
    else:
        return "Unsupported file type"

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python extract_text.py <file_path>")
        sys.exit(1)

    file_path = sys.argv[1]
    text = detect_and_extract_text(file_path)
    print(text.encode('utf-8'))
